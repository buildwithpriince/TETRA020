"""Document type detection + text extraction with page/slide/cell source refs."""
from __future__ import annotations

import csv
import io
import logging
import os
import re
from dataclasses import dataclass, field
from typing import List, Optional

log = logging.getLogger("prism.parsers")

# ---------- Detection --------------------------------------------------------

DETECTED_TYPES = {"pitch_deck", "mis", "financials", "projections", "cap_table", "unknown"}

# Keyword banks used for both filename hints and per-sheet content scoring.
_KEYWORDS = {
    "pitch_deck": ["pitch", "deck", "investor", "series a", "seed round", "vision",
                   "roadmap", "team", "market opportunity", "go-to-market", "traction"],
    "mis": ["mis", "monthly", "mtd", "management information", "mis dashboard",
            "monthly p&l", "monthly pnl", "monthly report", "weekly cohort"],
    "financials": ["financial", "p&l", "pnl", "profit and loss", "balance sheet",
                   "audited", "auditor", "notes to accounts", "trial balance",
                   "cash flow statement", "statutory"],
    "projections": ["projection", "forecast", "plan", "3yr", "5yr", "fy25", "fy26",
                    "fy27", "fy28", "base case", "bull case", "bear case",
                    "5-year plan", "3-year plan"],
    "cap_table": ["captable", "cap table", "cap_table", "cap-table", "ownership",
                  "shareholder", "shareholding", "shares", "esop pool", "waterfall",
                  "fully diluted"],
}
# Month names alone are a weak MIS signal (they also appear in projections).
_MONTHS = ["jan", "feb", "mar", "apr", "may", "jun", "jul", "aug", "sep", "oct", "nov", "dec"]


def _ext(filename: str) -> str:
    return os.path.splitext(filename)[1].lower().lstrip(".")


def _score_filename(filename: str) -> tuple[Optional[str], float]:
    """Score a filename against every candidate type; return (best, confidence)."""
    n = filename.lower().replace("-", " ").replace("_", " ")
    scores: dict[str, int] = {}
    for kind, kws in _KEYWORDS.items():
        for kw in kws:
            if kw in n:
                # Longer keywords beat short ones; count matches.
                scores[kind] = scores.get(kind, 0) + max(1, len(kw.split()))
    if not scores:
        return None, 0.0
    best = max(scores, key=scores.get)  # type: ignore[arg-type]
    # Confidence: 0.9 for unique strong hit, 0.75 for close race.
    top = scores[best]
    runner = sorted(scores.values(), reverse=True)[1] if len(scores) > 1 else 0
    conf = 0.9 if top >= 2 * max(runner, 1) else 0.75
    return best, conf


def _score_text(text: str) -> tuple[Optional[str], float, dict[str, int]]:
    """Score raw text against each candidate type. Returns (best, confidence, scores)."""
    n = text.lower()
    scores: dict[str, int] = {k: 0 for k in _KEYWORDS}
    for kind, kws in _KEYWORDS.items():
        for kw in kws:
            if kw in n:
                scores[kind] += 2 if len(kw.split()) > 1 else 1
    # Month hits nudge toward MIS but only weakly (projections often list months too).
    month_hits = sum(1 for m in _MONTHS if m in n)
    if month_hits >= 6:
        scores["mis"] += 3
    elif month_hits >= 3:
        scores["mis"] += 1

    total = sum(scores.values())
    if total == 0:
        return None, 0.0, scores
    best = max(scores, key=scores.get)  # type: ignore[arg-type]
    top = scores[best]
    conf = min(0.9, 0.4 + 0.1 * top)
    return best, conf, scores


def detect_type_from_filename(filename: str) -> Optional[str]:
    kind, _ = _score_filename(filename)
    return kind


def _sample_xlsx_text(path: str, max_sheets: int = 8, rows_per_sheet: int = 6) -> str:
    """Concatenate every sheet's title + its first N rows into one blob for scoring."""
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []
    for i, sheet in enumerate(wb.worksheets):
        if i >= max_sheets:
            break
        parts.append(f"[sheet:{sheet.title}]")
        for r_idx, row in enumerate(
            sheet.iter_rows(min_row=1, max_row=rows_per_sheet, values_only=True),
            start=1,
        ):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" ".join(cells))
    return "\n".join(parts)


def _sample_csv_text(path: str, max_rows: int = 20) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows: List[str] = []
        for i, row in enumerate(reader):
            if i >= max_rows:
                break
            rows.append(" ".join(row))
    return "\n".join(rows)


def detect_type_from_content(filename: str, path: str) -> tuple[str, float]:
    """Detect DocumentType by combining filename + content scoring.

    Strategy:
      1. Score the filename. If the filename gives a confident (>=0.9) match, trust it.
      2. Score the content. Merge with filename score. Best combined wins.
      3. Special cases: pptx→pitch_deck, PDF financials keywords→financials, cap-table
         detection dominates when ownership/shares signals appear alongside %.
    """
    ext = _ext(filename)
    file_kind, file_conf = _score_filename(filename)

    # Deck by extension (unless filename says otherwise strongly).
    if ext == "pptx":
        if file_kind and file_conf >= 0.9 and file_kind != "pitch_deck":
            return file_kind, file_conf
        return file_kind or "pitch_deck", max(file_conf, 0.95)

    if ext == "pdf":
        try:
            import pdfplumber

            with pdfplumber.open(path) as pdf:
                sample = "\n".join(
                    (pdf.pages[i].extract_text() or "") for i in range(min(3, len(pdf.pages)))
                )
            content_kind, content_conf, _ = _score_text(sample)
            return _combine(file_kind, file_conf, content_kind, content_conf, default="pitch_deck")
        except Exception as e:  # noqa: BLE001
            log.warning("pdf sniff failed for %s: %s", filename, e)
            return file_kind or "pitch_deck", max(file_conf, 0.55)

    if ext in ("xlsx", "xls"):
        try:
            sample = _sample_xlsx_text(path)
            content_kind, content_conf, scores = _score_text(sample)
            # Cap-table override: % + (shareholder/shares/ownership) dominant.
            if "%" in sample and scores.get("cap_table", 0) >= 2:
                return "cap_table", 0.9
            return _combine(file_kind, file_conf, content_kind, content_conf, default="mis")
        except Exception as e:  # noqa: BLE001
            log.warning("xls sniff failed for %s: %s", filename, e)
            return file_kind or "mis", max(file_conf, 0.55)

    if ext == "csv":
        try:
            sample = _sample_csv_text(path)
            content_kind, content_conf, scores = _score_text(sample)
            if "%" in sample and scores.get("cap_table", 0) >= 2:
                return "cap_table", 0.9
            return _combine(file_kind, file_conf, content_kind, content_conf, default="mis")
        except Exception as e:  # noqa: BLE001
            log.warning("csv sniff failed for %s: %s", filename, e)
            return file_kind or "mis", max(file_conf, 0.55)

    return file_kind or "unknown", max(file_conf, 0.4)


def _combine(
    file_kind: Optional[str],
    file_conf: float,
    content_kind: Optional[str],
    content_conf: float,
    default: str,
) -> tuple[str, float]:
    """Merge filename + content signals into a single decision.

    - If filename is confident (>=0.9), trust it.
    - Else if they agree, boost confidence.
    - Else pick whichever has higher confidence.
    """
    if file_kind and file_conf >= 0.9:
        return file_kind, file_conf
    if file_kind and content_kind and file_kind == content_kind:
        return file_kind, min(0.95, max(file_conf, content_conf) + 0.1)
    if content_kind and (content_conf > file_conf):
        return content_kind, content_conf
    if file_kind:
        return file_kind, file_conf
    if content_kind:
        return content_kind, content_conf
    return default, 0.5


# ---------- Extraction ------------------------------------------------------


@dataclass
class DocumentChunk:
    """A piece of text with its source reference."""

    text: str
    source_ref: str


@dataclass
class ParsedDocument:
    filename: str
    detected_type: str
    chunks: List[DocumentChunk] = field(default_factory=list)

    def joined(self, max_chars: int = 20000) -> str:
        buf: List[str] = []
        total = 0
        for c in self.chunks:
            block = f"[{c.source_ref}]\n{c.text}\n"
            if total + len(block) > max_chars:
                break
            buf.append(block)
            total += len(block)
        return "\n".join(buf)


def parse_document(path: str, filename: str, detected_type: str) -> ParsedDocument:
    ext = _ext(filename)
    if ext == "pdf":
        return _parse_pdf(path, filename, detected_type)
    if ext == "pptx":
        return _parse_pptx(path, filename, detected_type)
    if ext == "csv":
        return _parse_csv(path, filename, detected_type)
    if ext in ("xlsx", "xls"):
        return _parse_xlsx(path, filename, detected_type)
    return ParsedDocument(filename=filename, detected_type=detected_type)


def _parse_pdf(path: str, filename: str, detected_type: str) -> ParsedDocument:
    doc = ParsedDocument(filename=filename, detected_type=detected_type)
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                txt = (page.extract_text() or "").strip()
                if txt:
                    doc.chunks.append(DocumentChunk(text=txt, source_ref=f"{filename} p.{i}"))
    except Exception as e:  # noqa: BLE001
        log.warning("pdf parse fail %s: %s", filename, e)
    return doc


def _parse_pptx(path: str, filename: str, detected_type: str) -> ParsedDocument:
    doc = ParsedDocument(filename=filename, detected_type=detected_type)
    try:
        from pptx import Presentation

        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            if parts:
                doc.chunks.append(
                    DocumentChunk(text="\n".join(parts), source_ref=f"{filename} slide {i}")
                )
    except Exception as e:  # noqa: BLE001
        log.warning("pptx parse fail %s: %s", filename, e)
    return doc


def _parse_csv(path: str, filename: str, detected_type: str) -> ParsedDocument:
    doc = ParsedDocument(filename=filename, detected_type=detected_type)
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            buf = io.StringIO()
            for i, row in enumerate(reader, start=1):
                buf.write(f"row {i}: " + ", ".join(row) + "\n")
                if i > 300:
                    break
            doc.chunks.append(DocumentChunk(text=buf.getvalue(), source_ref=f"{filename} rows 1..N"))
    except Exception as e:  # noqa: BLE001
        log.warning("csv parse fail %s: %s", filename, e)
    return doc


def _parse_xlsx(path: str, filename: str, detected_type: str) -> ParsedDocument:
    doc = ParsedDocument(filename=filename, detected_type=detected_type)
    try:
        from openpyxl import load_workbook

        wb = load_workbook(path, data_only=True, read_only=True)
        for sheet in wb.worksheets:
            buf = io.StringIO()
            for r_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    buf.write(f"row {r_idx}: " + ", ".join(cells) + "\n")
                if r_idx > 250:
                    break
            content = buf.getvalue().strip()
            if content:
                doc.chunks.append(
                    DocumentChunk(text=content, source_ref=f'{filename} sheet "{sheet.title}"')
                )
    except Exception as e:  # noqa: BLE001
        log.warning("xlsx parse fail %s: %s", filename, e)
    return doc
