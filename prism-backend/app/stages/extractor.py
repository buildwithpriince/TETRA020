"""
Stage 2 -- extractor.py

Parses raw text/tables out of each validated file, per-file-type. Output is
a plain intermediate representation (page/slide/sheet -> text + tables with
a source_ref-friendly locator) that Stage 3-5 consume. No AI here -- this is
pure parsing.
"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field

logger = logging.getLogger("prism.extractor")


@dataclass
class ExtractedUnit:
    """One addressable chunk of a document -- a slide, a page, a sheet."""

    locator: str  # e.g. "slide 8", "page 3", "sheet 'MIS Jun-25', row 12"
    text: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)  # list of tables (rows of cells)


@dataclass
class ExtractedDocument:
    filename: str
    file_type_hint: str  # extension-based hint, e.g. "pdf", "pptx", "xlsx", "csv"
    units: list[ExtractedUnit] = field(default_factory=list)
    parse_error: str | None = None

    def full_text(self) -> str:
        return "\n\n".join(f"[{u.locator}]\n{u.text}" for u in self.units if u.text)


def _extract_pdf(data: bytes) -> list[ExtractedUnit]:
    import pdfplumber

    units: list[ExtractedUnit] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            tables = [t for t in (page.extract_tables() or []) if t]
            units.append(ExtractedUnit(locator=f"page {i}", text=text, tables=tables))
    return units


def _extract_pptx(data: bytes) -> list[ExtractedUnit]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    units: list[ExtractedUnit] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        tables: list[list[list[str]]] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t.strip():
                    texts.append(t)
            if shape.has_table:
                tbl = shape.table
                rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                tables.append(rows)
        # speaker notes can carry claims too (e.g. "assumes CAC stays at $50")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            texts.append("[notes] " + slide.notes_slide.notes_text_frame.text)
        units.append(ExtractedUnit(locator=f"slide {i}", text="\n".join(texts), tables=tables))
    return units


def _extract_xlsx(data: bytes) -> list[ExtractedUnit]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    units: list[ExtractedUnit] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append(["" if c is None else str(c) for c in row])
        # cap absurdly large sheets so we don't blow the prompt budget later
        rows = rows[:2000]
        units.append(ExtractedUnit(locator=f"sheet '{sheet_name}'", text="", tables=[rows] if rows else []))
    return units


def _extract_csv(data: bytes) -> list[ExtractedUnit]:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(data))
    rows = [df.columns.tolist()] + df.astype(str).values.tolist()
    return [ExtractedUnit(locator="csv", tables=[rows])]


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
    "csv": _extract_csv,
}


def extract(filename: str, data: bytes) -> ExtractedDocument:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        return ExtractedDocument(filename=filename, file_type_hint=ext, parse_error=f"Unsupported extension: {ext}")
    try:
        units = fn(data)
        return ExtractedDocument(filename=filename, file_type_hint=ext, units=units)
    except Exception as exc:
        logger.exception("Extraction failed for %s", filename)
        return ExtractedDocument(filename=filename, file_type_hint=ext, parse_error=str(exc))
